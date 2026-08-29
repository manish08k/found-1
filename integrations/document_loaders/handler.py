"""
Document loader nodes — load content from various sources for RAG pipelines.

Nodes:
  - loader.text            — plain text input
  - loader.json            — parse and flatten JSON
  - loader.csv             — CSV rows as documents
  - loader.pdf             — PDF text extraction (via pdfminer or pypdf2)
  - loader.web_scrape      — scrape a URL
  - loader.sitemap         — crawl a sitemap XML and scrape listed URLs
  - loader.github_repo     — fetch files from a GitHub repo
  - loader.youtube         — YouTube transcript via API
  - loader.confluence      — Confluence space/page loader
  - loader.s3_file         — load a file from S3
  - loader.notion_page     — load a Notion page (via Notion API)
"""
import asyncio
import csv
import io
import json
import re

import httpx
import structlog

from core.execution_engine import register_node
from core.config import settings

log = structlog.get_logger(__name__)


def _chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> list[str]:
    """Default chunking: fixed-size windows with overlap."""
    if len(text) <= chunk_size:
        return [text]
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += chunk_size - overlap
    return chunks


# ─── Plain Text ────────────────────────────────────────────────────────────────

@register_node("loader.text")
async def loader_text(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """Wraps a string as a document list ready for vector.upsert."""
    text = config.get("text") or input_data.get("text", "")
    metadata = config.get("metadata") or {}
    chunk_size = int(config.get("chunk_size", 0))

    if chunk_size > 0:
        chunks = _chunk_text(text, chunk_size, int(config.get("chunk_overlap", 200)))
        docs = [{"text": c, "metadata": {**metadata, "chunk_index": i}} for i, c in enumerate(chunks)]
    else:
        docs = [{"text": text, "metadata": metadata}]

    return {"documents": docs, "count": len(docs), "source": "text"}


# ─── JSON Loader ───────────────────────────────────────────────────────────────

@register_node("loader.json")
async def loader_json(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """
    Converts a JSON object or array into documents.
    - pointer: JSONPointer-style path to the array of items (e.g. "/items")
    - text_keys: list of fields whose values form the document text
    """
    raw = config.get("json") or input_data.get("json") or input_data
    if isinstance(raw, str):
        raw = json.loads(raw)

    pointer = config.get("pointer", "")
    if pointer:
        for part in pointer.strip("/").split("/"):
            if isinstance(raw, dict):
                raw = raw.get(part)
            elif isinstance(raw, list) and part.isdigit():
                raw = raw[int(part)]

    items = raw if isinstance(raw, list) else [raw]
    text_keys = config.get("text_keys") or []
    metadata_keys = config.get("metadata_keys") or []

    docs = []
    for item in items:
        if not isinstance(item, dict):
            text = str(item)
            meta = {}
        elif text_keys:
            text = " ".join(str(item.get(k, "")) for k in text_keys)
            meta = {k: item.get(k) for k in metadata_keys} if metadata_keys else {}
        else:
            text = json.dumps(item)
            meta = {}
        docs.append({"text": text, "metadata": meta})

    return {"documents": docs, "count": len(docs), "source": "json"}


# ─── CSV Loader ────────────────────────────────────────────────────────────────

@register_node("loader.csv")
async def loader_csv(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """
    Parses CSV content into documents (one per row).
    - content: the CSV string
    - text_columns: which columns to combine as document text
    """
    content = config.get("content") or input_data.get("content", "")
    text_columns = config.get("text_columns") or []
    delimiter = config.get("delimiter", ",")

    reader = csv.DictReader(io.StringIO(content), delimiter=delimiter)
    docs = []
    for i, row in enumerate(reader):
        if text_columns:
            text = " ".join(str(row.get(c, "")) for c in text_columns)
        else:
            text = " ".join(f"{k}: {v}" for k, v in row.items())
        docs.append({"text": text, "metadata": dict(row), "row_index": i})

    return {"documents": docs, "count": len(docs), "source": "csv"}


# ─── PDF Loader ────────────────────────────────────────────────────────────────

@register_node("loader.pdf")
async def loader_pdf(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """
    Extract text from a PDF (base64-encoded bytes or a URL).
    Requires pypdf: pip install pypdf
    """
    import base64

    url = config.get("url") or input_data.get("url")
    pdf_b64 = config.get("pdf_base64") or input_data.get("pdf_base64")
    chunk_size = int(config.get("chunk_size", 1000))
    chunk_overlap = int(config.get("chunk_overlap", 200))

    if url:
        from core.ssrf_guard import assert_safe_url
        assert_safe_url(url)
        async with httpx.AsyncClient(timeout=30) as client:
            r = await client.get(url)
            r.raise_for_status()
            pdf_bytes = r.content
    elif pdf_b64:
        pdf_bytes = base64.b64decode(pdf_b64)
    else:
        raise ValueError("loader.pdf requires 'url' or 'pdf_base64'")

    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(pdf_bytes))
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                pages.append({"text": text, "page": i + 1})
    except ImportError:
        raise RuntimeError("loader.pdf requires pypdf: pip install pypdf")

    docs = []
    for page_info in pages:
        chunks = _chunk_text(page_info["text"], chunk_size, chunk_overlap)
        for j, chunk in enumerate(chunks):
            docs.append({"text": chunk, "metadata": {"page": page_info["page"], "chunk_index": j}})

    return {"documents": docs, "count": len(docs), "pages": len(pages), "source": "pdf"}


# ─── Web Scraper ───────────────────────────────────────────────────────────────

@register_node("loader.web_scrape")
async def loader_web_scrape(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """
    Scrapes a URL and returns the page text.
    Requires html2text or beautifulsoup4 for clean extraction.
    """
    url = config.get("url") or input_data.get("url")
    if not url:
        raise ValueError("loader.web_scrape requires 'url'")

    from core.ssrf_guard import assert_safe_url
    assert_safe_url(url)

    selector = config.get("selector")  # optional CSS selector
    chunk_size = int(config.get("chunk_size", 1000))
    chunk_overlap = int(config.get("chunk_overlap", 200))

    async with httpx.AsyncClient(timeout=30, follow_redirects=True) as client:
        r = await client.get(url, headers={"User-Agent": "AutoFlow/1.0"})
        r.raise_for_status()
        html = r.text

    # Try BeautifulSoup for clean extraction
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(html, "html.parser")
        # Remove script/style tags
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        if selector:
            elements = soup.select(selector)
            text = " ".join(el.get_text(separator=" ", strip=True) for el in elements)
        else:
            text = soup.get_text(separator=" ", strip=True)
    except ImportError:
        # Fallback: strip HTML tags with regex
        text = re.sub(r"<[^>]+>", " ", html)
        text = re.sub(r"\s+", " ", text).strip()

    chunks = _chunk_text(text, chunk_size, chunk_overlap) if text else []
    docs = [{"text": c, "metadata": {"url": url, "chunk_index": i}} for i, c in enumerate(chunks)]

    return {
        "documents": docs,
        "count": len(docs),
        "url": url,
        "title": re.search(r"<title>(.*?)</title>", html, re.IGNORECASE | re.DOTALL),
        "source": "web_scrape",
    }


# ─── Sitemap Crawler ───────────────────────────────────────────────────────────

@register_node("loader.sitemap")
async def loader_sitemap(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """
    Reads a sitemap XML and scrapes each listed URL.
    """
    sitemap_url = config.get("sitemap_url") or input_data.get("sitemap_url")
    if not sitemap_url:
        raise ValueError("loader.sitemap requires 'sitemap_url'")

    from core.ssrf_guard import assert_safe_url
    assert_safe_url(sitemap_url)

    max_pages = int(config.get("max_pages", 20))
    chunk_size = int(config.get("chunk_size", 1000))

    async with httpx.AsyncClient(timeout=30) as client:
        r = await client.get(sitemap_url, headers={"User-Agent": "AutoFlow/1.0"})
        r.raise_for_status()
        xml = r.text

    urls = re.findall(r"<loc>(.*?)</loc>", xml)[:max_pages]
    all_docs = []

    async def scrape_one(url):
        try:
            assert_safe_url(url)
            async with httpx.AsyncClient(timeout=15, follow_redirects=True) as c:
                resp = await c.get(url, headers={"User-Agent": "AutoFlow/1.0"})
                resp.raise_for_status()
                html = resp.text
            try:
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(html, "html.parser")
                for tag in soup(["script", "style"]):
                    tag.decompose()
                text = soup.get_text(separator=" ", strip=True)
            except ImportError:
                text = re.sub(r"<[^>]+>", " ", html)
                text = re.sub(r"\s+", " ", text).strip()
            chunks = _chunk_text(text, chunk_size)
            return [{"text": c, "metadata": {"url": url, "chunk_index": i}} for i, c in enumerate(chunks)]
        except Exception as e:
            log.warning("sitemap_scrape_failed", url=url, error=str(e))
            return []

    results = await asyncio.gather(*[scrape_one(u) for u in urls])
    for page_docs in results:
        all_docs.extend(page_docs)

    return {"documents": all_docs, "count": len(all_docs), "urls_scraped": len(urls), "source": "sitemap"}


# ─── GitHub Repo Loader ────────────────────────────────────────────────────────

@register_node("loader.github_repo")
async def loader_github_repo(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """
    Loads text files from a GitHub repository via the GitHub API.
    """
    repo = config.get("repo") or input_data.get("repo")  # e.g. "owner/repo"
    if not repo:
        raise ValueError("loader.github_repo requires 'repo' (e.g. 'owner/name')")

    branch = config.get("branch", "main")
    file_patterns = config.get("file_patterns", ["*.py", "*.md", "*.txt"])
    max_files = int(config.get("max_files", 50))
    token = getattr(settings, "GITHUB_TOKEN", None)

    headers = {"Accept": "application/vnd.github+json"}
    if token:
        headers["Authorization"] = f"Bearer {token}"

    async with httpx.AsyncClient(timeout=30) as client:
        # Get the tree
        r = await client.get(
            f"https://api.github.com/repos/{repo}/git/trees/{branch}?recursive=1",
            headers=headers,
        )
        r.raise_for_status()
        tree = r.json()

    files = [f for f in tree.get("tree", []) if f.get("type") == "blob"]
    # Filter by extension pattern
    import fnmatch
    filtered = []
    for f in files:
        path = f.get("path", "")
        for pat in file_patterns:
            if fnmatch.fnmatch(path, pat):
                filtered.append(f)
                break

    filtered = filtered[:max_files]

    async def load_file(file_info):
        async with httpx.AsyncClient(timeout=15) as c:
            r = await c.get(
                f"https://api.github.com/repos/{repo}/contents/{file_info['path']}?ref={branch}",
                headers=headers,
            )
            if r.status_code != 200:
                return None
            data = r.json()
            import base64
            content = base64.b64decode(data.get("content", "")).decode("utf-8", errors="replace")
            return {"text": content, "metadata": {"path": file_info["path"], "repo": repo, "sha": file_info.get("sha")}}

    docs_raw = await asyncio.gather(*[load_file(f) for f in filtered])
    docs = [d for d in docs_raw if d]

    return {"documents": docs, "count": len(docs), "repo": repo, "branch": branch, "source": "github_repo"}


# ─── YouTube Transcript ────────────────────────────────────────────────────────

@register_node("loader.youtube")
async def loader_youtube(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """
    Fetches YouTube video transcript.
    Requires youtube-transcript-api: pip install youtube-transcript-api
    """
    video_id = config.get("video_id") or input_data.get("video_id")
    url = config.get("url") or input_data.get("url")

    if not video_id and url:
        # Extract video ID from URL
        m = re.search(r"(?:v=|youtu\.be/)([A-Za-z0-9_-]{11})", url)
        video_id = m.group(1) if m else None

    if not video_id:
        raise ValueError("loader.youtube requires 'video_id' or 'url'")

    try:
        from youtube_transcript_api import YouTubeTranscriptApi
        loop = asyncio.get_event_loop()
        transcript = await loop.run_in_executor(
            None,
            lambda: YouTubeTranscriptApi.get_transcript(video_id, languages=config.get("languages", ["en"]))
        )
    except ImportError:
        raise RuntimeError("loader.youtube requires youtube-transcript-api: pip install youtube-transcript-api")

    full_text = " ".join(item["text"] for item in transcript)
    chunk_size = int(config.get("chunk_size", 1000))
    chunks = _chunk_text(full_text, chunk_size)
    docs = [{"text": c, "metadata": {"video_id": video_id, "chunk_index": i}} for i, c in enumerate(chunks)]

    return {"documents": docs, "count": len(docs), "video_id": video_id, "source": "youtube"}


# ─── Notion Page Loader ────────────────────────────────────────────────────────

@register_node("loader.notion_page")
async def loader_notion_page(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """
    Loads a Notion page's blocks as text.
    Requires a Notion integration token (NOTION_API_KEY or credential_id).
    """
    page_id = config.get("page_id") or input_data.get("page_id")
    if not page_id:
        raise ValueError("loader.notion_page requires 'page_id'")

    api_key = getattr(settings, "NOTION_API_KEY", None)
    if not api_key and credential_id:
        # Decrypt credential
        from credentials.encryption import decrypt_token
        api_key = await decrypt_token(credential_id, db)

    if not api_key:
        raise ValueError("loader.notion_page requires NOTION_API_KEY or credential_id")

    headers = {
        "Authorization": f"Bearer {api_key}",
        "Notion-Version": "2022-06-28",
    }

    async def get_blocks(block_id):
        async with httpx.AsyncClient(timeout=30) as client:
            r = await client.get(
                f"https://api.notion.com/v1/blocks/{block_id}/children",
                headers=headers,
            )
            r.raise_for_status()
            return r.json().get("results", [])

    def extract_text(block):
        btype = block.get("type", "")
        content = block.get(btype, {})
        rich_text = content.get("rich_text", [])
        return "".join(rt.get("plain_text", "") for rt in rich_text)

    blocks = await get_blocks(page_id)
    lines = [extract_text(b) for b in blocks if extract_text(b)]
    text = "\n".join(lines)

    chunk_size = int(config.get("chunk_size", 1000))
    chunks = _chunk_text(text, chunk_size)
    docs = [{"text": c, "metadata": {"page_id": page_id, "chunk_index": i}} for i, c in enumerate(chunks)]

    return {"documents": docs, "count": len(docs), "page_id": page_id, "source": "notion"}


# ─── S3 File Loader ────────────────────────────────────────────────────────────

@register_node("loader.s3_file")
async def loader_s3_file(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """
    Load and parse a file from S3.
    Supports .txt, .csv, .json, .pdf (with pypdf).
    """
    try:
        import boto3
    except ImportError:
        raise RuntimeError("loader.s3_file requires boto3: pip install boto3")

    bucket = config.get("bucket") or input_data.get("bucket")
    key = config.get("key") or input_data.get("key")  # S3 object key

    if not bucket or not key:
        raise ValueError("loader.s3_file requires 'bucket' and 'key'")

    access_key = getattr(settings, "AWS_ACCESS_KEY_ID", None)
    secret_key = getattr(settings, "AWS_SECRET_ACCESS_KEY", None)
    region = config.get("region") or getattr(settings, "AWS_REGION", "us-east-1")

    s3 = boto3.client(
        "s3",
        region_name=region,
        aws_access_key_id=access_key,
        aws_secret_access_key=secret_key,
    )

    loop = asyncio.get_event_loop()
    obj = await loop.run_in_executor(None, lambda: s3.get_object(Bucket=bucket, Key=key))
    content_bytes = obj["Body"].read()

    ext = key.rsplit(".", 1)[-1].lower()
    chunk_size = int(config.get("chunk_size", 1000))

    if ext == "pdf":
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(content_bytes))
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
    elif ext == "csv":
        content_str = content_bytes.decode("utf-8", errors="replace")
        reader = csv.reader(io.StringIO(content_str))
        text = "\n".join(",".join(row) for row in reader)
    elif ext == "json":
        parsed = json.loads(content_bytes)
        text = json.dumps(parsed, indent=2)
    else:
        text = content_bytes.decode("utf-8", errors="replace")

    chunks = _chunk_text(text, chunk_size)
    docs = [{"text": c, "metadata": {"bucket": bucket, "key": key, "chunk_index": i}} for i, c in enumerate(chunks)]

    return {"documents": docs, "count": len(docs), "bucket": bucket, "key": key, "source": "s3"}


# ─── DOCX Loader ───────────────────────────────────────────────────────────────

@register_node("loader.docx")
async def loader_docx(config: dict, input_data: dict, credential_id: str | None, db) -> dict:
    """
    Load text from a .docx file (Microsoft Word).
    Requires python-docx: pip install python-docx
    config: file_path (local path) OR file_bytes (base64-encoded bytes), metadata
    """
    import base64

    file_path = config.get("file_path") or input_data.get("file_path")
    file_b64 = config.get("file_bytes") or input_data.get("file_bytes")
    metadata = config.get("metadata") or {}

    try:
        import docx as _docx
    except ImportError:
        raise RuntimeError("loader.docx requires python-docx: pip install python-docx")

    import io

    if file_b64:
        data = io.BytesIO(base64.b64decode(file_b64))
        doc = _docx.Document(data)
    elif file_path:
        doc = _docx.Document(file_path)
    else:
        raise ValueError("loader.docx requires 'file_path' or 'file_bytes'")

    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    full_text = "\n\n".join(paragraphs)

    chunk_size = int(config.get("chunk_size", 0))
    if chunk_size > 0:
        chunks = _chunk_text(full_text, chunk_size, int(config.get("chunk_overlap", 200)))
        docs = [{"text": c, "metadata": {**metadata, "chunk_index": i}} for i, c in enumerate(chunks)]
    else:
        docs = [{"text": full_text, "metadata": {**metadata, "source": file_path or "docx"}}]

    return {"documents": docs, "count": len(docs), "source": "docx"}


# ─── Folder Loader ─────────────────────────────────────────────────────────────

@register_node("loader.folder")
async def loader_folder(config: dict, input_data: dict, credential_id: str | None, db) -> dict:
    """
    Load all text files from a local folder recursively.
    config: folder_path, glob_pattern (default **/*.txt), extensions (list), metadata, recursive
    """
    import glob as _glob
    import os

    folder_path = config.get("folder_path") or input_data.get("folder_path")
    if not folder_path:
        raise ValueError("loader.folder requires 'folder_path'")

    extensions = config.get("extensions") or [".txt", ".md", ".csv", ".json"]
    pattern = config.get("glob_pattern", "**/*")
    recursive = config.get("recursive", True)
    metadata_base = config.get("metadata") or {}
    chunk_size = int(config.get("chunk_size", 0))

    files = _glob.glob(os.path.join(folder_path, pattern), recursive=recursive)
    files = [f for f in files if os.path.isfile(f) and any(f.endswith(ext) for ext in extensions)]

    docs = []
    for filepath in sorted(files):
        try:
            with open(filepath, encoding="utf-8", errors="ignore") as fh:
                text = fh.read()
            if not text.strip():
                continue
            rel_path = os.path.relpath(filepath, folder_path)
            meta = {**metadata_base, "source": filepath, "filename": os.path.basename(filepath), "path": rel_path}
            if chunk_size > 0:
                for i, chunk in enumerate(_chunk_text(text, chunk_size, int(config.get("chunk_overlap", 200)))):
                    docs.append({"text": chunk, "metadata": {**meta, "chunk_index": i}})
            else:
                docs.append({"text": text, "metadata": meta})
        except Exception as e:
            log.warning("folder_loader_file_error", filepath=filepath, error=str(e))

    return {"documents": docs, "count": len(docs), "files_processed": len(files), "source": "folder"}


# ─── Google Drive Loader ───────────────────────────────────────────────────────

@register_node("loader.google_drive")
async def loader_google_drive(config: dict, input_data: dict, credential_id: str | None, db) -> dict:
    """
    Load documents from Google Drive (files in a folder or by ID).
    Requires GOOGLE service account credentials or OAuth access token.
    config: access_token, folder_id (OR file_id), mime_types, max_files
    """
    access_token = config.get("access_token") or getattr(settings, "GOOGLE_ACCESS_TOKEN", "")
    if not access_token:
        raise ValueError("loader.google_drive requires 'access_token' or GOOGLE_ACCESS_TOKEN")

    headers = {"Authorization": f"Bearer {access_token}"}
    folder_id = config.get("folder_id") or input_data.get("folder_id")
    file_id = config.get("file_id") or input_data.get("file_id")
    max_files = int(config.get("max_files", 20))
    metadata_base = config.get("metadata") or {}
    docs = []

    async with httpx.AsyncClient(timeout=30) as client:
        if file_id:
            # Single file
            meta_r = await client.get(
                f"https://www.googleapis.com/drive/v3/files/{file_id}",
                params={"fields": "id,name,mimeType"},
                headers=headers,
            )
            meta_r.raise_for_status()
            file_meta = meta_r.json()
            files_list = [file_meta]
        elif folder_id:
            list_r = await client.get(
                "https://www.googleapis.com/drive/v3/files",
                params={"q": f"'{folder_id}' in parents and trashed=false",
                        "fields": "files(id,name,mimeType)", "pageSize": max_files},
                headers=headers,
            )
            list_r.raise_for_status()
            files_list = list_r.json().get("files", [])
        else:
            raise ValueError("loader.google_drive requires 'folder_id' or 'file_id'")

        for f in files_list[:max_files]:
            fid = f["id"]
            fname = f["name"]
            mime = f.get("mimeType", "")

            # Export Google Docs as plain text; download binary files
            if "google-apps.document" in mime:
                dl_r = await client.get(
                    f"https://www.googleapis.com/drive/v3/files/{fid}/export",
                    params={"mimeType": "text/plain"},
                    headers=headers,
                )
            else:
                dl_r = await client.get(
                    f"https://www.googleapis.com/drive/v3/files/{fid}",
                    params={"alt": "media"},
                    headers=headers,
                )

            if dl_r.status_code == 200:
                text = dl_r.text
                docs.append({"text": text, "metadata": {**metadata_base, "source": fid, "filename": fname, "mime_type": mime}})

    return {"documents": docs, "count": len(docs), "source": "google_drive"}


# ─── Google Sheets Loader ──────────────────────────────────────────────────────

@register_node("loader.google_sheets")
async def loader_google_sheets(config: dict, input_data: dict, credential_id: str | None, db) -> dict:
    """
    Load data from a Google Sheets spreadsheet as documents (one doc per row).
    config: access_token, spreadsheet_id, sheet_name (default first sheet), header_row (bool)
    """
    access_token = config.get("access_token") or getattr(settings, "GOOGLE_ACCESS_TOKEN", "")
    if not access_token:
        raise ValueError("loader.google_sheets requires 'access_token' or GOOGLE_ACCESS_TOKEN")

    spreadsheet_id = config.get("spreadsheet_id") or input_data.get("spreadsheet_id")
    if not spreadsheet_id:
        raise ValueError("loader.google_sheets requires 'spreadsheet_id'")

    sheet_name = config.get("sheet_name", "")
    range_str = f"{sheet_name}!A:ZZ" if sheet_name else "A:ZZ"
    has_header = config.get("header_row", True)
    metadata_base = config.get("metadata") or {}

    async with httpx.AsyncClient(timeout=30) as client:
        r = await client.get(
            f"https://sheets.googleapis.com/v4/spreadsheets/{spreadsheet_id}/values/{range_str}",
            headers={"Authorization": f"Bearer {access_token}"},
        )
        r.raise_for_status()
        data = r.json()

    values = data.get("values", [])
    if not values:
        return {"documents": [], "count": 0, "source": "google_sheets"}

    headers = values[0] if has_header else [f"col_{i}" for i in range(len(values[0]))]
    rows = values[1:] if has_header else values

    docs = []
    for i, row in enumerate(rows):
        row_dict = dict(zip(headers, row + [""] * (len(headers) - len(row))))
        text = " | ".join(f"{k}: {v}" for k, v in row_dict.items() if v)
        docs.append({
            "text": text,
            "metadata": {**metadata_base, "row": i + 2, "source": spreadsheet_id, "data": row_dict},
        })

    return {"documents": docs, "count": len(docs), "source": "google_sheets"}


# ─── Jira Loader ──────────────────────────────────────────────────────────────

@register_node("loader.jira")
async def loader_jira(config: dict, input_data: dict, credential_id: str | None, db) -> dict:
    """
    Load Jira issues as documents using JQL.
    config: base_url, email, api_token, jql (query), max_results, fields
    """
    import base64

    base_url = config.get("base_url") or getattr(settings, "JIRA_BASE_URL", "")
    email = config.get("email") or getattr(settings, "JIRA_EMAIL", "")
    api_token = config.get("api_token") or getattr(settings, "JIRA_API_TOKEN", "")

    if not all([base_url, email, api_token]):
        raise ValueError("loader.jira requires base_url, email, and api_token")

    jql = config.get("jql") or input_data.get("jql", "ORDER BY created DESC")
    max_results = int(config.get("max_results", 50))
    fields = config.get("fields") or ["summary", "description", "status", "assignee", "priority", "created"]

    auth = base64.b64encode(f"{email}:{api_token}".encode()).decode()
    headers = {"Authorization": f"Basic {auth}", "Accept": "application/json"}

    async with httpx.AsyncClient(timeout=30) as client:
        r = await client.get(
            f"{base_url.rstrip('/')}/rest/api/3/search",
            params={"jql": jql, "maxResults": max_results, "fields": ",".join(fields)},
            headers=headers,
        )
        r.raise_for_status()
        data = r.json()

    docs = []
    for issue in data.get("issues", []):
        f = issue.get("fields", {})
        summary = f.get("summary", "")
        description = f.get("description") or ""
        if isinstance(description, dict):
            # Jira ADF format — extract plain text
            description = " ".join(
                c.get("text", "") for block in description.get("content", [])
                for c in block.get("content", []) if isinstance(c, dict)
            )
        text = f"{summary}\n\n{description}".strip()
        docs.append({
            "text": text,
            "metadata": {
                "source": issue.get("key"),
                "url": f"{base_url}/browse/{issue.get('key')}",
                "status": f.get("status", {}).get("name"),
                "priority": f.get("priority", {}).get("name"),
                "assignee": (f.get("assignee") or {}).get("displayName"),
            },
        })

    return {"documents": docs, "count": len(docs), "source": "jira"}


# ─── FireCrawl Loader ─────────────────────────────────────────────────────────

@register_node("loader.firecrawl")
async def loader_firecrawl(config: dict, input_data: dict, credential_id: str | None, db) -> dict:
    """
    Crawl and scrape URLs using FireCrawl API (JS rendering, clean markdown output).
    Requires FIRECRAWL_API_KEY.
    config: url, mode (scrape | crawl), max_pages (crawl only), formats
    """
    api_key = getattr(settings, "FIRECRAWL_API_KEY", None)
    if not api_key:
        raise ValueError("loader.firecrawl requires FIRECRAWL_API_KEY")

    url = config.get("url") or input_data.get("url")
    if not url:
        raise ValueError("loader.firecrawl requires 'url'")

    mode = config.get("mode", "scrape")
    metadata_base = config.get("metadata") or {}
    formats = config.get("formats") or ["markdown"]

    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    docs = []

    async with httpx.AsyncClient(timeout=60) as client:
        if mode == "scrape":
            r = await client.post(
                "https://api.firecrawl.dev/v1/scrape",
                json={"url": url, "formats": formats},
                headers=headers,
            )
            r.raise_for_status()
            data = r.json().get("data", {})
            text = data.get("markdown") or data.get("content") or data.get("text", "")
            docs.append({"text": text, "metadata": {**metadata_base, "source": url,
                                                     "title": data.get("metadata", {}).get("title", "")}})
        else:
            # Crawl mode
            max_pages = int(config.get("max_pages", 10))
            r = await client.post(
                "https://api.firecrawl.dev/v1/crawl",
                json={"url": url, "limit": max_pages, "scrapeOptions": {"formats": formats}},
                headers=headers,
            )
            r.raise_for_status()
            crawl_id = r.json().get("id")

            # Poll for completion
            import asyncio as _asyncio
            for _ in range(30):
                await _asyncio.sleep(3)
                poll_r = await client.get(f"https://api.firecrawl.dev/v1/crawl/{crawl_id}", headers=headers)
                poll_r.raise_for_status()
                status_data = poll_r.json()
                if status_data.get("status") == "completed":
                    for page in status_data.get("data", []):
                        text = page.get("markdown") or page.get("content", "")
                        docs.append({"text": text, "metadata": {**metadata_base,
                                                                  "source": page.get("metadata", {}).get("sourceURL", url),
                                                                  "title": page.get("metadata", {}).get("title", "")}})
                    break

    return {"documents": docs, "count": len(docs), "source": "firecrawl"}


# ─── JSON Lines Loader ─────────────────────────────────────────────────────────

@register_node("loader.jsonlines")
async def loader_jsonlines(config: dict, input_data: dict, credential_id: str | None, db) -> dict:
    """
    Load JSONL (newline-delimited JSON) as documents.
    config: content (JSONL string) OR file_path, text_key, metadata_keys, max_lines
    """
    content = config.get("content") or input_data.get("content", "")
    file_path = config.get("file_path") or input_data.get("file_path")
    text_key = config.get("text_key", "text")
    metadata_keys = config.get("metadata_keys") or []
    max_lines = int(config.get("max_lines", 10000))
    metadata_base = config.get("metadata") or {}

    if file_path and not content:
        with open(file_path, encoding="utf-8") as fh:
            content = fh.read()

    if not content:
        raise ValueError("loader.jsonlines requires 'content' or 'file_path'")

    import json as _json
    docs = []
    for i, line in enumerate(content.strip().splitlines()):
        if i >= max_lines:
            break
        line = line.strip()
        if not line:
            continue
        try:
            obj = _json.loads(line)
        except Exception:
            obj = {"text": line}

        text = obj.get(text_key) or " ".join(str(v) for v in obj.values())
        meta = {**metadata_base, "line": i + 1}
        for k in metadata_keys:
            if k in obj:
                meta[k] = obj[k]
        docs.append({"text": str(text), "metadata": meta})

    return {"documents": docs, "count": len(docs), "source": "jsonlines"}


# ─── EPUB Loader ──────────────────────────────────────────────────────────────

@register_node("loader.epub")
async def loader_epub(config: dict, input_data: dict, credential_id: str | None, db) -> dict:
    """
    Load text from an EPUB ebook file.
    Requires ebooklib + beautifulsoup4: pip install EbookLib beautifulsoup4
    config: file_path OR file_bytes (base64), chunk_size, metadata
    """
    import base64
    import io

    file_path = config.get("file_path") or input_data.get("file_path")
    file_b64 = config.get("file_bytes") or input_data.get("file_bytes")
    metadata_base = config.get("metadata") or {}
    chunk_size = int(config.get("chunk_size", 0))

    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
    except ImportError:
        raise RuntimeError("loader.epub requires EbookLib and beautifulsoup4: pip install EbookLib beautifulsoup4")

    if file_b64:
        book = epub.read_epub(io.BytesIO(base64.b64decode(file_b64)))
    elif file_path:
        book = epub.read_epub(file_path)
    else:
        raise ValueError("loader.epub requires 'file_path' or 'file_bytes'")

    chapters = []
    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            soup = BeautifulSoup(item.get_content(), "html.parser")
            text = soup.get_text(separator="\n", strip=True)
            if text.strip():
                chapters.append(text)

    full_text = "\n\n".join(chapters)

    if chunk_size > 0:
        chunks = _chunk_text(full_text, chunk_size, int(config.get("chunk_overlap", 200)))
        docs = [{"text": c, "metadata": {**metadata_base, "chunk_index": i}} for i, c in enumerate(chunks)]
    else:
        docs = [{"text": c, "metadata": {**metadata_base, "chapter": i}} for i, c in enumerate(chapters)]

    return {"documents": docs, "count": len(docs), "source": "epub"}


# ─── GitBook Loader ───────────────────────────────────────────────────────────

@register_node("loader.gitbook")
async def loader_gitbook(config: dict, input_data: dict, credential_id: str | None, db) -> dict:
    """
    Load pages from a GitBook space via GitBook API v1.
    Requires GITBOOK_API_TOKEN.
    config: space_id, api_token (overrides env), max_pages, metadata
    """
    api_token = config.get("api_token") or getattr(settings, "GITBOOK_API_TOKEN", None)
    if not api_token:
        raise ValueError("loader.gitbook requires GITBOOK_API_TOKEN or 'api_token'")

    space_id = config.get("space_id") or input_data.get("space_id")
    if not space_id:
        raise ValueError("loader.gitbook requires 'space_id'")

    max_pages = int(config.get("max_pages", 100))
    metadata_base = config.get("metadata") or {}
    headers = {"Authorization": f"Bearer {api_token}", "Accept": "application/json"}

    async with httpx.AsyncClient(timeout=30) as client:
        # List pages
        r = await client.get(
            f"https://api.gitbook.com/v1/spaces/{space_id}/content",
            headers=headers,
        )
        r.raise_for_status()
        data = r.json()

    pages = data.get("pages", [])[:max_pages]
    docs = []

    async with httpx.AsyncClient(timeout=30) as client:
        for page in pages:
            page_id = page.get("id")
            title = page.get("title", "")
            path = page.get("path", "")
            try:
                pr = await client.get(
                    f"https://api.gitbook.com/v1/spaces/{space_id}/content/page/{page_id}",
                    headers={**headers, "Accept": "text/plain"},
                )
                text = pr.text if pr.status_code == 200 else title
            except Exception:
                text = title

            if text.strip():
                docs.append({"text": text, "metadata": {**metadata_base,
                                                         "source": f"{space_id}/{page_id}",
                                                         "title": title, "path": path}})

    return {"documents": docs, "count": len(docs), "source": "gitbook"}


# ─── API Document Loader ───────────────────────────────────────────────────────

@register_node("loader.api")
async def loader_api(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """Load documents from a REST API endpoint."""
    from core.ssrf_guard import assert_safe_url, SSRFSafeTransport
    url = config.get("url") or input_data.get("url")
    if not url:
        raise ValueError("loader.api requires 'url'")
    assert_safe_url(url)
    method = config.get("method", "GET").upper()
    headers = dict(config.get("headers") or {})
    body = config.get("body") or config.get("params")
    data_path = config.get("data_path", "")
    text_field = config.get("text_field", "")
    metadata_fields = config.get("metadata_fields") or []
    chunk_size = int(config.get("chunk_size", 0))
    async with httpx.AsyncClient(timeout=30, transport=SSRFSafeTransport()) as client:
        if method == "POST":
            r = await client.post(url, json=body, headers=headers)
        else:
            r = await client.get(url, params=body, headers=headers)
        r.raise_for_status()
    try:
        data = r.json()
    except Exception:
        data = r.text
    # Navigate to data_path
    if data_path and isinstance(data, dict):
        for key in data_path.split("."):
            data = data.get(key, data) if isinstance(data, dict) else data
    items = data if isinstance(data, list) else [data]
    documents = []
    for item in items:
        if isinstance(item, dict):
            text = str(item.get(text_field, "") if text_field else item)
            meta = {f: item.get(f) for f in metadata_fields if f in item}
        else:
            text = str(item)
            meta = {}
        meta["source"] = url
        if chunk_size > 0:
            chunks = _chunk_text(text, chunk_size)
            for i, c in enumerate(chunks):
                documents.append({"text": c, "metadata": {**meta, "chunk_index": i}})
        else:
            documents.append({"text": text, "metadata": meta})
    return {"documents": documents, "total": len(documents), "source": url}


# ─── Apify Website Content Crawler ────────────────────────────────────────────

@register_node("loader.apify_crawler")
async def loader_apify_crawler(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """Load documents by crawling websites via Apify."""
    api_key = getattr(settings, "APIFY_API_KEY", None)
    if not api_key:
        raise ValueError("loader.apify_crawler requires APIFY_API_KEY in environment")
    url = config.get("url") or input_data.get("url")
    if not url:
        raise ValueError("loader.apify_crawler requires 'url'")
    max_pages = int(config.get("max_pages", 10))
    chunk_size = int(config.get("chunk_size", 1000))
    # Start the crawl run
    async with httpx.AsyncClient(timeout=300) as client:
        run_resp = await client.post(
            "https://api.apify.com/v2/acts/apify~website-content-crawler/runs",
            headers={"Authorization": f"Bearer {api_key}"},
            json={"startUrls": [{"url": url}], "maxCrawlingDepth": 1, "maxPagesPerCrawl": max_pages},
        )
        run_resp.raise_for_status()
        run_id = run_resp.json()["data"]["id"]
        # Poll for completion
        for _ in range(60):
            await asyncio.sleep(5)
            status_resp = await client.get(
                f"https://api.apify.com/v2/acts/apify~website-content-crawler/runs/{run_id}",
                headers={"Authorization": f"Bearer {api_key}"},
            )
            status_resp.raise_for_status()
            status = status_resp.json()["data"]["status"]
            if status in ("SUCCEEDED", "FAILED", "TIMED-OUT", "ABORTED"):
                break
        # Fetch dataset items
        dataset_id = status_resp.json()["data"]["defaultDatasetId"]
        items_resp = await client.get(
            f"https://api.apify.com/v2/datasets/{dataset_id}/items",
            headers={"Authorization": f"Bearer {api_key}"},
            params={"format": "json", "clean": "true"},
        )
        items_resp.raise_for_status()
        items = items_resp.json()
    documents = []
    for item in items:
        text = item.get("text") or item.get("markdown") or item.get("html") or str(item)
        meta = {"url": item.get("url", ""), "title": item.get("title", ""), "source": "apify"}
        if chunk_size > 0:
            for i, c in enumerate(_chunk_text(text, chunk_size)):
                documents.append({"text": c, "metadata": {**meta, "chunk_index": i}})
        else:
            documents.append({"text": text, "metadata": meta})
    return {"documents": documents, "total": len(documents), "source": url}


# ─── Brave Search as Document Loader ──────────────────────────────────────────

@register_node("loader.brave_search_docs")
async def loader_brave_search_docs(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """Load search results from Brave Search as documents."""
    api_key = getattr(settings, "BRAVE_SEARCH_API_KEY", None)
    if not api_key:
        raise ValueError("loader.brave_search_docs requires BRAVE_SEARCH_API_KEY")
    query = config.get("query") or input_data.get("query", "")
    if not query:
        raise ValueError("loader.brave_search_docs requires 'query'")
    count = min(int(config.get("count", 10)), 20)
    async with httpx.AsyncClient(timeout=30) as client:
        r = await client.get(
            "https://api.search.brave.com/res/v1/web/search",
            headers={"Accept": "application/json", "X-Subscription-Token": api_key},
            params={"q": query, "count": count},
        )
        r.raise_for_status()
        data = r.json()
    results = data.get("web", {}).get("results", [])
    documents = []
    for item in results:
        text = f"{item.get('title', '')}\n{item.get('description', '')}"
        documents.append({
            "text": text,
            "metadata": {"url": item.get("url", ""), "title": item.get("title", ""), "source": "brave_search"},
        })
    return {"documents": documents, "total": len(documents), "query": query}


# ─── Cheerio Web Scraper ──────────────────────────────────────────────────────

@register_node("loader.cheerio")
async def loader_cheerio(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """Scrape a webpage and extract content by CSS selector (Python equivalent)."""
    from core.ssrf_guard import assert_safe_url, SSRFSafeTransport
    url = config.get("url") or input_data.get("url")
    if not url:
        raise ValueError("loader.cheerio requires 'url'")
    assert_safe_url(url)
    selector = config.get("selector", "body")
    include_links = config.get("include_links", False)
    chunk_size = int(config.get("chunk_size", 1000))
    async with httpx.AsyncClient(timeout=30, transport=SSRFSafeTransport(),
                                  headers={"User-Agent": "Mozilla/5.0 AutoFlow/1.0"}) as client:
        r = await client.get(url)
        r.raise_for_status()
        html = r.text
    # Parse with html.parser
    try:
        from html.parser import HTMLParser
        import re as _re

        class TextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.texts = []
                self.links = []
                self._skip = False
                self._skip_tags = {"script", "style", "head", "meta", "link"}
                self._current_tag = None

            def handle_starttag(self, tag, attrs):
                self._current_tag = tag
                if tag in self._skip_tags:
                    self._skip = True
                if include_links and tag == "a":
                    href = dict(attrs).get("href", "")
                    if href:
                        self.links.append(href)

            def handle_endtag(self, tag):
                if tag in self._skip_tags:
                    self._skip = False

            def handle_data(self, data):
                if not self._skip and data.strip():
                    self.texts.append(data.strip())

        parser = TextExtractor()
        # Simple tag extraction - extract text within the "selector" tag
        tag_match = _re.search(rf"<{selector.lstrip('#. ')}[^>]*>(.*?)</{selector.lstrip('#. ')}>",
                               html, _re.DOTALL | _re.IGNORECASE)
        target_html = tag_match.group(1) if tag_match else html
        parser.feed(target_html)
        text = " ".join(parser.texts)
        if include_links:
            text += "\n\nLinks: " + ", ".join(parser.links)
    except Exception:
        # Fallback: strip all tags
        import re as _re
        text = _re.sub(r"<[^>]+>", " ", html)
        text = _re.sub(r"\s+", " ", text).strip()

    meta = {"url": url, "selector": selector, "source": "cheerio"}
    documents = []
    if chunk_size > 0:
        for i, c in enumerate(_chunk_text(text, chunk_size)):
            documents.append({"text": c, "metadata": {**meta, "chunk_index": i}})
    else:
        documents.append({"text": text, "metadata": meta})
    return {"documents": documents, "total": len(documents), "source": url}


# ─── Custom Document Loader ───────────────────────────────────────────────────

@register_node("loader.custom_document")
async def loader_custom_document(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """Pass-through loader: accepts documents directly as configuration."""
    documents = (
        config.get("documents")
        or input_data.get("documents")
        or []
    )
    # Normalize each document
    normalized = []
    for doc in documents:
        if isinstance(doc, str):
            normalized.append({"text": doc, "metadata": {}})
        elif isinstance(doc, dict):
            text = doc.get("text") or doc.get("content") or str(doc)
            meta = doc.get("metadata") or {}
            normalized.append({"text": text, "metadata": meta})
    return {"documents": normalized, "total": len(normalized), "source": "custom"}


# ─── Document Store Loader ────────────────────────────────────────────────────

@register_node("loader.document_store")
async def loader_document_store(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """Load documents from the platform's internal document store."""
    store_id = config.get("store_id") or input_data.get("store_id")
    limit = int(config.get("limit", 100))
    try:
        from sqlalchemy import select
        from storage.models import DocumentStore
        q = select(DocumentStore)
        if store_id:
            q = q.where(DocumentStore.id == store_id)
        q = q.limit(limit)
        result = await db.execute(q)
        stores = result.scalars().all()
        documents = []
        for store in stores:
            docs = store.documents if hasattr(store, "documents") else []
            for doc in (docs if isinstance(docs, list) else []):
                text = doc.get("text") or doc.get("content", "")
                meta = doc.get("metadata") or {"store_id": str(store.id)}
                documents.append({"text": text, "metadata": meta})
        return {"documents": documents, "total": len(documents), "store_id": store_id}
    except Exception as e:
        log.warning("loader_document_store_error", error=str(e))
        return {"documents": [], "total": 0, "store_id": store_id, "error": str(e)}


# ─── Figma File Loader ────────────────────────────────────────────────────────

@register_node("loader.figma")
async def loader_figma(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """Load design data from a Figma file."""
    api_key = getattr(settings, "FIGMA_ACCESS_TOKEN", None)
    if not api_key:
        raise ValueError("loader.figma requires FIGMA_ACCESS_TOKEN in environment")
    file_key = config.get("file_key") or input_data.get("file_key")
    if not file_key:
        raise ValueError("loader.figma requires 'file_key'")
    node_ids = config.get("node_ids")
    async with httpx.AsyncClient(timeout=60) as client:
        params = {}
        if node_ids:
            params["ids"] = ",".join(node_ids) if isinstance(node_ids, list) else node_ids
        r = await client.get(
            f"https://api.figma.com/v1/files/{file_key}",
            headers={"X-Figma-Token": api_key},
            params=params,
        )
        r.raise_for_status()
        data = r.json()
    documents = []

    def extract_nodes(node, depth=0):
        if depth > 5:
            return
        name = node.get("name", "")
        node_type = node.get("type", "")
        text = ""
        if node_type == "TEXT":
            text = node.get("characters", "")
        if text or node_type in ("FRAME", "COMPONENT", "INSTANCE"):
            desc = f"[{node_type}] {name}"
            if text:
                desc += f": {text}"
            documents.append({
                "text": desc,
                "metadata": {"type": node_type, "name": name, "file_key": file_key, "source": "figma"},
            })
        for child in node.get("children", []):
            extract_nodes(child, depth + 1)

    doc = data.get("document", {})
    extract_nodes(doc)
    return {"documents": documents, "total": len(documents), "file_key": file_key}


# ─── Generic File Loader ──────────────────────────────────────────────────────

@register_node("loader.file")
async def loader_file(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """Load a file by path or base64 content, dispatching to appropriate sub-loader."""
    import base64, os, tempfile
    from core.execution_engine import NODE_HANDLERS
    file_path = config.get("file_path") or input_data.get("file_path")
    file_content_b64 = config.get("file_content") or input_data.get("file_content")
    file_type = (config.get("file_type") or "").lower()

    if file_path:
        ext = os.path.splitext(file_path)[1].lower().lstrip(".")
        file_type = file_type or ext
        with open(file_path, "rb") as f:
            raw = f.read()
    elif file_content_b64:
        raw = base64.b64decode(file_content_b64)
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_type}") as tmp:
            tmp.write(raw)
            file_path = tmp.name
    else:
        raise ValueError("loader.file requires 'file_path' or 'file_content' (base64)")

    type_map = {
        "pdf": "loader.pdf", "txt": "loader.text", "csv": "loader.csv",
        "json": "loader.json", "jsonl": "loader.jsonlines",
        "docx": "loader.docx", "doc": "loader.docx",
        "xlsx": "loader.microsoft_excel", "xls": "loader.microsoft_excel",
        "pptx": "loader.microsoft_powerpoint", "epub": "loader.epub",
    }
    sub_node = type_map.get(file_type)
    if sub_node and sub_node in NODE_HANDLERS:
        sub_config = {**config, "file_path": file_path}
        return await NODE_HANDLERS[sub_node](sub_config, input_data, credential_id, db)

    # Fallback: plain text
    try:
        text = raw.decode("utf-8", errors="replace")
    except Exception:
        text = str(raw[:10000])
    return {"documents": [{"text": text, "metadata": {"source": file_path}}], "total": 1}


# ─── Microsoft Excel Loader ───────────────────────────────────────────────────

@register_node("loader.microsoft_excel")
async def loader_microsoft_excel(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """Load Excel (.xlsx) files as documents — each row becomes a document."""
    import base64, io, tempfile, os
    file_path = config.get("file_path") or input_data.get("file_path")
    file_content_b64 = config.get("file_content") or input_data.get("file_content")
    sheet_name = config.get("sheet_name")  # None = all sheets
    row_as_text = config.get("row_as_text", True)

    try:
        import openpyxl
    except ImportError:
        raise ImportError("loader.microsoft_excel requires openpyxl: pip install openpyxl")

    if file_content_b64:
        raw = base64.b64decode(file_content_b64)
        wb = openpyxl.load_workbook(io.BytesIO(raw), data_only=True)
    elif file_path:
        wb = openpyxl.load_workbook(file_path, data_only=True)
    else:
        raise ValueError("loader.microsoft_excel requires 'file_path' or 'file_content'")

    sheets = [wb[sheet_name]] if sheet_name else wb.worksheets
    documents = []
    for ws in sheets:
        headers = [str(cell.value or "") for cell in next(ws.iter_rows(min_row=1, max_row=1))]
        for row in ws.iter_rows(min_row=2, values_only=True):
            row_dict = dict(zip(headers, [str(v) if v is not None else "" for v in row]))
            text = " | ".join(f"{k}: {v}" for k, v in row_dict.items() if v) if row_as_text else str(row_dict)
            documents.append({"text": text, "metadata": {"sheet": ws.title, "row": row_dict, "source": file_path or "excel"}})
    return {"documents": documents, "total": len(documents), "sheets": [ws.title for ws in sheets]}


# ─── Microsoft PowerPoint Loader ──────────────────────────────────────────────

@register_node("loader.microsoft_powerpoint")
async def loader_microsoft_powerpoint(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """Load PowerPoint (.pptx) files — each slide becomes a document."""
    import base64, io
    file_path = config.get("file_path") or input_data.get("file_path")
    file_content_b64 = config.get("file_content") or input_data.get("file_content")
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("loader.microsoft_powerpoint requires python-pptx: pip install python-pptx")

    if file_content_b64:
        raw = base64.b64decode(file_content_b64)
        prs = Presentation(io.BytesIO(raw))
    elif file_path:
        prs = Presentation(file_path)
    else:
        raise ValueError("loader.microsoft_powerpoint requires 'file_path' or 'file_content'")

    documents = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        text = "\n".join(texts)
        if text:
            documents.append({
                "text": text,
                "metadata": {"slide": i + 1, "source": file_path or "pptx"},
            })
    return {"documents": documents, "total": len(documents), "slides": len(prs.slides)}


# ─── Microsoft Word Loader ────────────────────────────────────────────────────

@register_node("loader.microsoft_word")
async def loader_microsoft_word(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """Load Word (.docx) files as documents."""
    # Alias to loader.docx with enhanced options
    from core.execution_engine import NODE_HANDLERS
    if "loader.docx" in NODE_HANDLERS:
        return await NODE_HANDLERS["loader.docx"](config, input_data, credential_id, db)
    # Fallback implementation
    import base64, io
    file_path = config.get("file_path") or input_data.get("file_path")
    file_content_b64 = config.get("file_content") or input_data.get("file_content")
    try:
        import docx
    except ImportError:
        raise ImportError("loader.microsoft_word requires python-docx: pip install python-docx")
    if file_content_b64:
        raw = base64.b64decode(file_content_b64)
        doc = docx.Document(io.BytesIO(raw))
    elif file_path:
        doc = docx.Document(file_path)
    else:
        raise ValueError("loader.microsoft_word requires 'file_path' or 'file_content'")
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return {"documents": [{"text": text, "metadata": {"source": file_path or "docx"}}], "total": 1}


# ─── Oxylabs Web Scraper Loader ───────────────────────────────────────────────

@register_node("loader.oxylabs")
async def loader_oxylabs(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """Load web content via Oxylabs Realtime API."""
    username = getattr(settings, "OXYLABS_USERNAME", None)
    password = getattr(settings, "OXYLABS_PASSWORD", None)
    if not username or not password:
        raise ValueError("loader.oxylabs requires OXYLABS_USERNAME and OXYLABS_PASSWORD")
    url = config.get("url") or input_data.get("url")
    if not url:
        raise ValueError("loader.oxylabs requires 'url'")
    parse = config.get("parse", True)
    async with httpx.AsyncClient(timeout=60) as client:
        r = await client.post(
            "https://realtime.oxylabs.io/v1/queries",
            auth=(username, password),
            json={"source": "universal", "url": url, "parse": parse},
        )
        r.raise_for_status()
        data = r.json()
    results = data.get("results", [])
    documents = []
    for result in results:
        content = result.get("content") or result.get("parsed", {}).get("text", "") or str(result)
        documents.append({"text": content, "metadata": {"url": url, "source": "oxylabs"}})
    return {"documents": documents, "total": len(documents), "source": url}


# ─── Playwright Web Scraper ───────────────────────────────────────────────────

@register_node("loader.playwright_web")
async def loader_playwright_web(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """Scrape JS-heavy sites using Playwright (or httpx fallback)."""
    from core.ssrf_guard import assert_safe_url, SSRFSafeTransport
    url = config.get("url") or input_data.get("url")
    if not url:
        raise ValueError("loader.playwright_web requires 'url'")
    assert_safe_url(url)
    chunk_size = int(config.get("chunk_size", 1000))
    import re as _re
    try:
        from playwright.async_api import async_playwright
        async with async_playwright() as p:
            browser = await p.chromium.launch(headless=True)
            page = await browser.new_page()
            await page.goto(url, timeout=30000)
            if config.get("wait_for_selector"):
                await page.wait_for_selector(config["wait_for_selector"], timeout=10000)
            content = await page.content()
            await browser.close()
        text = _re.sub(r"<[^>]+>", " ", content)
        text = _re.sub(r"\s+", " ", text).strip()
    except ImportError:
        # Fallback to httpx
        async with httpx.AsyncClient(timeout=30, transport=SSRFSafeTransport(),
                                     headers={"User-Agent": "Mozilla/5.0"}) as client:
            r = await client.get(url)
            text = _re.sub(r"<[^>]+>", " ", r.text)
            text = _re.sub(r"\s+", " ", text).strip()
    meta = {"url": url, "source": "playwright"}
    documents = []
    if chunk_size > 0:
        for i, c in enumerate(_chunk_text(text, chunk_size)):
            documents.append({"text": c, "metadata": {**meta, "chunk_index": i}})
    else:
        documents.append({"text": text, "metadata": meta})
    return {"documents": documents, "total": len(documents), "source": url}


# ─── Puppeteer-style Web Scraper ──────────────────────────────────────────────

@register_node("loader.puppeteer_web")
async def loader_puppeteer_web(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """Scrape sites with Playwright (Puppeteer equivalent in Python)."""
    # Delegate to playwright loader since Puppeteer is Node.js only
    from core.execution_engine import NODE_HANDLERS
    if "loader.playwright_web" in NODE_HANDLERS:
        return await NODE_HANDLERS["loader.playwright_web"](config, input_data, credential_id, db)
    raise RuntimeError("loader.puppeteer_web: requires playwright — install with: pip install playwright && playwright install chromium")


# ─── S3 Directory Loader ──────────────────────────────────────────────────────

@register_node("loader.s3_directory")
async def loader_s3_directory(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """Load all files from an S3 prefix as documents."""
    bucket = config.get("bucket") or input_data.get("bucket")
    prefix = config.get("prefix", "")
    file_extensions = config.get("file_extensions") or [".txt", ".pdf", ".csv", ".json", ".md"]
    max_files = int(config.get("max_files", 50))
    if not bucket:
        raise ValueError("loader.s3_directory requires 'bucket'")
    try:
        import boto3
        s3 = boto3.client(
            "s3",
            aws_access_key_id=settings.AWS_ACCESS_KEY_ID or None,
            aws_secret_access_key=settings.AWS_SECRET_ACCESS_KEY or None,
            region_name=settings.AWS_REGION,
        )
        response = s3.list_objects_v2(Bucket=bucket, Prefix=prefix)
        objects = response.get("Contents", [])[:max_files]
        from core.execution_engine import NODE_HANDLERS
        documents = []
        for obj in objects:
            key = obj["Key"]
            ext = "." + key.rsplit(".", 1)[-1].lower() if "." in key else ""
            if file_extensions and ext not in file_extensions:
                continue
            body = s3.get_object(Bucket=bucket, Key=key)["Body"].read()
            try:
                text = body.decode("utf-8", errors="replace")
            except Exception:
                continue
            documents.append({"text": text, "metadata": {"s3_key": key, "bucket": bucket, "source": f"s3://{bucket}/{key}"}})
        return {"documents": documents, "total": len(documents), "bucket": bucket, "prefix": prefix}
    except ImportError:
        raise ImportError("loader.s3_directory requires boto3: pip install boto3")


# ─── SearchAPI Document Loader ────────────────────────────────────────────────

@register_node("loader.searchapi_docs")
async def loader_searchapi_docs(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """Load search results from SearchAPI.io as documents."""
    api_key = getattr(settings, "SEARCHAPI_API_KEY", None)
    if not api_key:
        raise ValueError("loader.searchapi_docs requires SEARCHAPI_API_KEY")
    query = config.get("query") or input_data.get("query", "")
    if not query:
        raise ValueError("loader.searchapi_docs requires 'query'")
    engine = config.get("engine", "google")
    num = int(config.get("num", 10))
    async with httpx.AsyncClient(timeout=30) as client:
        r = await client.get(
            "https://www.searchapi.io/api/v1/search",
            params={"api_key": api_key, "engine": engine, "q": query, "num": num},
        )
        r.raise_for_status()
        data = r.json()
    results = data.get("organic_results", [])
    documents = []
    for item in results:
        text = f"{item.get('title', '')}\n{item.get('snippet', '')}"
        documents.append({
            "text": text,
            "metadata": {"url": item.get("link", ""), "title": item.get("title", ""), "source": "searchapi"},
        })
    return {"documents": documents, "total": len(documents), "query": query}


# ─── SerpAPI Document Loader ──────────────────────────────────────────────────

@register_node("loader.serpapi_docs")
async def loader_serpapi_docs(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """Load search results from SerpAPI as documents."""
    api_key = getattr(settings, "SERPAPI_API_KEY", None)
    if not api_key:
        raise ValueError("loader.serpapi_docs requires SERPAPI_API_KEY")
    query = config.get("query") or input_data.get("query", "")
    if not query:
        raise ValueError("loader.serpapi_docs requires 'query'")
    engine = config.get("engine", "google")
    num_results = int(config.get("num_results", 10))
    async with httpx.AsyncClient(timeout=30) as client:
        r = await client.get(
            "https://serpapi.com/search",
            params={"api_key": api_key, "engine": engine, "q": query, "num": num_results},
        )
        r.raise_for_status()
        data = r.json()
    organic = data.get("organic_results", [])
    documents = []
    for item in organic:
        text = f"{item.get('title', '')}\n{item.get('snippet', '')}"
        documents.append({
            "text": text,
            "metadata": {"url": item.get("link", ""), "source": "serpapi"},
        })
    return {"documents": documents, "total": len(documents), "query": query}


# ─── Spider.cloud Crawler ─────────────────────────────────────────────────────

@register_node("loader.spider_crawler")
async def loader_spider_crawler(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """Crawl websites using Spider.cloud API."""
    api_key = getattr(settings, "SPIDER_API_KEY", None)
    if not api_key:
        raise ValueError("loader.spider_crawler requires SPIDER_API_KEY")
    url = config.get("url") or input_data.get("url")
    if not url:
        raise ValueError("loader.spider_crawler requires 'url'")
    limit = int(config.get("limit", 10))
    chunk_size = int(config.get("chunk_size", 1000))
    async with httpx.AsyncClient(timeout=120) as client:
        r = await client.post(
            "https://api.spider.cloud/crawl",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={"url": url, "limit": limit, "return_format": "markdown"},
        )
        r.raise_for_status()
        data = r.json()
    documents = []
    items = data if isinstance(data, list) else data.get("data", [])
    for item in items:
        text = item.get("markdown") or item.get("content") or str(item)
        meta = {"url": item.get("url", url), "source": "spider"}
        if chunk_size > 0:
            for i, c in enumerate(_chunk_text(text, chunk_size)):
                documents.append({"text": c, "metadata": {**meta, "chunk_index": i}})
        else:
            documents.append({"text": text, "metadata": meta})
    return {"documents": documents, "total": len(documents), "source": url}


# ─── Unstructured.io Loader ───────────────────────────────────────────────────

@register_node("loader.unstructured")
async def loader_unstructured(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """Parse documents via Unstructured.io API (PDF, DOCX, HTML, etc.)."""
    api_key = getattr(settings, "UNSTRUCTURED_API_KEY", None)
    if not api_key:
        raise ValueError("loader.unstructured requires UNSTRUCTURED_API_KEY")
    import base64, tempfile
    file_path = config.get("file_path") or input_data.get("file_path")
    file_content_b64 = config.get("file_content") or input_data.get("file_content")
    file_type = config.get("file_type", "pdf")
    strategy = config.get("strategy", "auto")
    if file_content_b64:
        raw = base64.b64decode(file_content_b64)
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_type}") as tmp:
            tmp.write(raw)
            file_path = tmp.name
    if not file_path:
        raise ValueError("loader.unstructured requires 'file_path' or 'file_content'")
    async with httpx.AsyncClient(timeout=120) as client:
        with open(file_path, "rb") as f:
            r = await client.post(
                "https://api.unstructured.io/general/v0/general",
                headers={"unstructured-api-key": api_key},
                files={"files": (file_path, f, f"application/{file_type}")},
                data={"strategy": strategy},
            )
            r.raise_for_status()
            elements = r.json()
    documents = []
    for el in elements:
        text = el.get("text", "")
        if text:
            documents.append({
                "text": text,
                "metadata": {
                    "type": el.get("type", ""),
                    "page_number": el.get("metadata", {}).get("page_number"),
                    "source": file_path,
                },
            })
    return {"documents": documents, "total": len(documents), "source": file_path}


# ─── VectorStore to Document Loader ──────────────────────────────────────────

@register_node("loader.vectorstore_to_doc")
async def loader_vectorstore_to_doc(config: dict, input_data: dict, credential_id: str, db) -> dict:
    """Convert vector store search results into documents for further processing."""
    from core.execution_engine import NODE_HANDLERS
    collection = config.get("collection", "default")
    vs_type = config.get("vectorstore_type", "inmemory")
    query = config.get("query") or input_data.get("query") or input_data.get("input", "")
    top_k = int(config.get("top_k", 10))
    if not query:
        raise ValueError("loader.vectorstore_to_doc requires 'query'")
    query_node = f"vectorstore.{vs_type}.query"
    if query_node not in NODE_HANDLERS:
        raise ValueError(f"loader.vectorstore_to_doc: vector store type '{vs_type}' not found")
    result = await NODE_HANDLERS[query_node](
        {"collection": collection, "query": query, "top_k": top_k},
        {"query": query},
        credential_id,
        db,
    )
    docs = result.get("results", result.get("documents", []))
    # Normalize to document format
    documents = []
    for doc in docs:
        text = doc.get("content") or doc.get("text", str(doc))
        meta = doc.get("metadata", {})
        meta["score"] = doc.get("score")
        meta["source"] = f"vectorstore.{vs_type}/{collection}"
        documents.append({"text": text, "metadata": meta})
    return {"documents": documents, "total": len(documents), "query": query, "collection": collection}
